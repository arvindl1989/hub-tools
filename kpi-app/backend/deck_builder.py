"""Marketing Hub monthly-sharing deck generator.

Two independent pieces:
  1. A generic {{token}} merge engine (`fill_pptx_template`) that knows nothing
     about KONE or ticket data — it just replaces `{{name}}` placeholders found
     inside run text anywhere in a .pptx (including tables and grouped shapes)
     with values from a dict. Pointing it at a different templatized deck later
     is just a new token map, no code change here.
  2. KONE-specific logic (`compute_marketing_deck_tokens`) that turns a ticket
     dataframe + the feedback sheet into the token values for
     templates/marketing_hub_deck.pptx specifically.
"""
import re
from io import BytesIO
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

TEMPLATE_PATH = Path(__file__).parent / "templates" / "marketing_hub_deck.pptx"

_TOKEN_RE = re.compile(r"\{\{(\w+)\}\}")

# ── Generic merge engine ────────────────────────────────────────────────────


def _iter_all_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_all_shapes(shape.shapes)
        else:
            yield shape


def _replace_tokens_in_text_frame(tf, tokens: dict) -> None:
    for para in tf.paragraphs:
        for run in para.runs:
            if "{{" in run.text:
                run.text = _TOKEN_RE.sub(lambda m: str(tokens.get(m.group(1), "")), run.text)


def fill_pptx_template(template_path: Path, tokens: dict, bar_widths: dict | None = None) -> bytes:
    """`bar_widths` is an optional {slide_index: {shape_id: new_width_emu}} map
    for shapes that represent a proportional bar (left edge stays put, width is
    set directly) — used for e.g. a leaderboard whose bar lengths must reflect
    real values rather than whatever the template's sample data happened to
    draw. MUST be keyed by slide index — shape IDs are only unique within a
    single slide, not across the presentation, so an unscoped {shape_id: width}
    map will silently clobber unrelated shapes on other slides that happen to
    reuse the same ID."""
    prs = Presentation(str(template_path))
    for slide_index, slide in enumerate(prs.slides):
        slide_bar_widths = (bar_widths or {}).get(slide_index, {})
        for shape in _iter_all_shapes(slide.shapes):
            if shape.has_text_frame:
                _replace_tokens_in_text_frame(shape.text_frame, tokens)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _replace_tokens_in_text_frame(cell.text_frame, tokens)
            if shape.shape_id in slide_bar_widths:
                shape.width = slide_bar_widths[shape.shape_id]
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── KONE Marketing Hub token computation ────────────────────────────────────

COMPLETED_STATES = {"Closed Completed", "Confirmation Completed"}

# Fixed roster for the "Feedback highlights" slide. Pooja is excluded from
# the leaderboard entirely (never appears, at any rank) but stays in the
# quote-bubble identity list — those are unrelated, six fixed name slots.
ROSTER = [
    {"key": "pooja",      "full_name": "Pooja V",                    "short_name": "Pooja",      "in_leaderboard": False, "has_quote_slot": False},
    {"key": "akshayaar",  "full_name": "Akshayaa Rajeswari AS",      "short_name": "Akshayaa R", "in_leaderboard": True,  "has_quote_slot": True},
    {"key": "arvind",     "full_name": "Arvind Lakshminarayanan",    "short_name": "Arvind",     "in_leaderboard": True,  "has_quote_slot": True},
    {"key": "ajith",      "full_name": "Ajith A",                    "short_name": "Ajith",      "in_leaderboard": True,  "has_quote_slot": True},
    {"key": "akshayap",   "full_name": "Akshaya Praveen",            "short_name": "Akshaya P",  "in_leaderboard": True,  "has_quote_slot": True},
    {"key": "ranjith",    "full_name": "Ranjithkumar Ashokkumar",    "short_name": "Ranjith",    "in_leaderboard": True,  "has_quote_slot": True},
    {"key": "nitish",     "full_name": "Nitish JK",                  "short_name": "Nitish",     "in_leaderboard": True,  "has_quote_slot": True},
]

# Slide indices (0-based) that carry a proportional-bar section. All bar
# geometry below (shape IDs, track widths) is only valid WITHIN its slide —
# shape IDs are not unique across the whole presentation — so every fill-id
# map here must stay paired with its slide index wherever it's used (see
# fill_pptx_template's bar_widths param).
KPI_SLIDE_INDEX = 2        # "Hub KPIs"
REQUESTS_SLIDE_INDEX = 3   # "Requests overview"
LEADERBOARD_SLIDE_INDEX = 7  # "Feedback highlights"

# Feedback Rating bars (slide 3) — one track/fill pair per rating parameter.
# Scale is absolute (score ÷ scale_max), not relative to each other — a
# perfect score is always a full bar regardless of what the other parameters
# scored.
RATING_BAR_FILL_IDS = {"overall": 16, "quality": 20, "timeliness": 24, "interaction": 28}
RATING_BAR_TRACK_WIDTH_EMU = 4533900

# Ticket Inflow by Area bars (slide 4) — one bar per area, scaled relative to
# the largest area count shown that period.
AREA_BAR_FILL_IDS = {"eu": 12, "apm": 16, "global": 20, "ame": 24, "gcn": 28}
AREA_BAR_TRACK_WIDTH_EMU = 2190750

# By Work Type bars (slide 4) — inflow+outflow pair per sub-category. All 10
# bars (5 categories × inflow/outflow) share ONE max so their lengths stay
# comparable to each other, matching the "Inflow vs Outflow" framing.
WORK_TYPE_BAR_FILL_IDS = {
    "wt_wcm_inflow": 37,   "wt_wcm_outflow": 40,
    "wt_cpgd_inflow": 44,  "wt_cpgd_outflow": 47,
    "wt_email_inflow": 51, "wt_email_outflow": 54,
    "wt_dcg_inflow": 58,   "wt_dcg_outflow": 61,
    "wt_ra_inflow": 65,    "wt_ra_outflow": 68,
}
WORK_TYPE_BAR_TRACK_WIDTH_EMU = 4048125

# "Feedbacks by team member" leaderboard (slide 8) — 6 physical row slots,
# each a name + count text plus a track (background) and fill (value) bar
# shape. Rank is assigned at generation time by sorting counts descending,
# so the bar for the highest count is always slot 1's full-width fill and
# every other bar is scaled relative to it.
LEADERBOARD_FILL_SHAPE_IDS = [25, 29, 33, 37, 41, 45]
LEADERBOARD_TRACK_WIDTH_EMU = 2628900

# area value (normalized upper) -> key-requests slide column key
AREA_TO_KEY_REQUEST_COL = {
    "EU": "europe", "APM": "apm", "AME": "ame", "GLOBAL": "global", "GCN": "global",
}
KEY_REQUEST_SLOTS = {"europe": 6, "apm": 5, "ame": 5, "global": 2}

# sub_category (matched by keyword, case-insensitive) -> work-type token prefix
WORK_TYPE_KEYWORDS = [
    ("wt_wcm",  "website content management"),
    ("wt_cpgd", "content production"),
    ("wt_email", "email"),
    ("wt_dcg",  "demand creation"),
    ("wt_ra",   "retention"),
]


def _bar_widths_by_max(values: dict, track_width_emu: int) -> dict:
    """{key: raw_count} -> {key: width_emu}, each bar scaled relative to the
    largest value in the group (that value's bar is full width)."""
    max_v = max(values.values(), default=0)
    if not max_v:
        return {k: 0 for k in values}
    return {k: round(track_width_emu * (v / max_v)) for k, v in values.items()}


def _ordinal(n: int) -> str:
    if 10 <= n % 100 <= 20:
        suf = "th"
    else:
        suf = {1: "st", 2: "nd", 3: "rd"}.get(n % 10, "th")
    return f"{n}{suf}"


def _period_labels(date_from: str, date_to: str) -> dict:
    frm = pd.Timestamp(date_from)
    to = pd.Timestamp(date_to)

    is_full_month = (
        frm.day == 1
        and to.normalize() == (frm + pd.offsets.MonthEnd(0)).normalize()
        and frm.year == to.year
        and frm.month == to.month
    )

    if is_full_month:
        month_year = frm.strftime("%B %Y")
        month_name = frm.strftime("%B")
    else:
        month_year = f"{_ordinal(frm.day)} {frm.strftime('%b')} – {_ordinal(to.day)} {to.strftime('%b %Y')}"
        month_name = "the selected period"

    return {
        "deck_month_year": month_year,
        "deck_month_year_caps": month_year.upper(),
        "deck_month_name": month_name,
        "period_from_fmt": f"{_ordinal(frm.day)} {frm.strftime('%B')}",
        "period_to_fmt": f"{_ordinal(to.day)} {to.strftime('%B %Y')}",
    }


def _fmt_num(v) -> str:
    return "" if v is None else str(int(v)) if float(v).is_integer() else f"{v:g}"


def _fmt_score(v, decimals=1) -> str:
    return "—" if v is None else f"{v:.{decimals}f}"


def compute_marketing_deck_tokens(
    df: pd.DataFrame,
    hub_health: dict,
    feedback_data: dict,
    date_from: str,
    date_to: str,
    generated_date: str,
) -> tuple[dict, dict]:
    """`df` is the raw ticket dataframe for the session (unfiltered by date —
    filtering happens here for the inflow/outflow/key-request breakdowns).
    `hub_health` is the dict returned by calling the /hub-health endpoint logic
    for the SAME date range — reused rather than recomputed so the deck's KPI
    numbers always match what the Dashboard itself shows. `feedback_data` is
    the dict returned by feedback_summary() for the same range. `generated_date`
    is a pre-formatted "D MONTH YYYY" string (deck build time, not the
    reporting period). Returns (tokens, bar_widths) — bar_widths is the
    {slide_index: {shape_id: width_emu}} map fill_pptx_template needs for the
    leaderboard bars."""
    tokens: dict = {"deck_generated_date": generated_date}
    tokens.update(_period_labels(date_from, date_to))

    created = df["created_date"] if "created_date" in df.columns else pd.Series(dtype="datetime64[ns]")
    closed = df["closed_date"] if "closed_date" in df.columns else pd.Series(dtype="datetime64[ns]")
    frm_ts, to_ts = pd.Timestamp(date_from), pd.Timestamp(date_to) + pd.Timedelta(days=1)

    inflow_mask = created.notna() & (created >= frm_ts) & (created < to_ts)
    outflow_mask = closed.notna() & (closed >= frm_ts) & (closed < to_ts)
    inflow_df = df[inflow_mask]
    outflow_df = df[outflow_mask]

    tokens["inflow_total"] = _fmt_num(len(inflow_df))
    tokens["outflow_total"] = _fmt_num(len(outflow_df))

    # Hub KPIs (slide 3) — same figures as the Dashboard's Hub Health / KPI tiles
    tokens["kpi_total_tickets"] = _fmt_num(hub_health.get("total", 0))
    tokens["kpi_resolved"] = _fmt_num(hub_health.get("closed_completed", 0))
    tokens["kpi_in_pipeline"] = _fmt_num(hub_health.get("in_pipeline", 0))

    # By-area inflow (slide 4)
    area_values = {col: 0 for col in ("eu", "apm", "ame", "global", "gcn")}
    if "area" in inflow_df.columns:
        area_counts = inflow_df.dropna(subset=["area"]).groupby("area").size()
        for area_val, count in area_counts.items():
            norm = str(area_val).strip().upper()
            if norm in {"EU", "APM", "AME", "GLOBAL", "GCN"}:
                area_values[norm.lower()] = int(count)
    for col, count in area_values.items():
        tokens[f"inflow_area_{col}"] = _fmt_num(count)

    # By work-type inflow/outflow (slide 4)
    work_type_values = {f"{prefix}_{d}": 0 for prefix, _ in WORK_TYPE_KEYWORDS for d in ("inflow", "outflow")}
    if "sub_category" in df.columns:
        for prefix, kw in WORK_TYPE_KEYWORDS:
            match = df["sub_category"].fillna("").str.lower().str.contains(kw, regex=False)
            work_type_values[f"{prefix}_inflow"] = int((inflow_mask & match).sum())
            work_type_values[f"{prefix}_outflow"] = int((outflow_mask & match).sum())
    for key, count in work_type_values.items():
        tokens[key] = _fmt_num(count)

    # Key requests (slide 7) — completed tickets in range, grouped by area,
    # most-recently-closed first, capped to each column's slot count.
    for col, n in KEY_REQUEST_SLOTS.items():
        for i in range(1, n + 1):
            tokens[f"key_request_{col}_{i}"] = ""
    if {"area", "state", "short_description"}.issubset(df.columns):
        delivered = outflow_df[outflow_df["state"].isin(COMPLETED_STATES)].dropna(subset=["area"])
        delivered = delivered.sort_values("closed_date", ascending=False)
        buckets: dict[str, list] = {"europe": [], "apm": [], "ame": [], "global": []}
        for _, row in delivered.iterrows():
            col = AREA_TO_KEY_REQUEST_COL.get(str(row["area"]).strip().upper())
            if not col:
                continue
            desc = str(row.get("short_description", "")).strip()
            if desc and len(buckets[col]) < KEY_REQUEST_SLOTS[col]:
                buckets[col].append(desc)
        for col, items in buckets.items():
            for i, desc in enumerate(items, start=1):
                tokens[f"key_request_{col}_{i}"] = desc

    # Feedback (slide 3 + 8)
    param_avgs = feedback_data.get("param_avgs") or {}
    scale_max = feedback_data.get("scale_max") or 5
    rating_scores = {
        "overall": param_avgs.get("overall", feedback_data.get("avg_score")),
        "quality": param_avgs.get("quality"),
        "timeliness": param_avgs.get("timeliness"),
        "interaction": param_avgs.get("interaction"),
    }
    tokens["kpi_feedbacks"] = _fmt_num(feedback_data.get("total") or 0)
    tokens["kpi_feedback_rating"] = _fmt_score(feedback_data.get("avg_score"), 2)
    tokens["kpi_overall"] = _fmt_score(rating_scores["overall"])
    tokens["kpi_quality"] = _fmt_score(rating_scores["quality"])
    tokens["kpi_timeliness"] = _fmt_score(rating_scores["timeliness"])
    tokens["kpi_interaction"] = _fmt_score(rating_scores["interaction"])

    by_user_count = {row["user"]: row["count"] for row in feedback_data.get("by_user", [])}
    entries = feedback_data.get("entries", [])  # already newest-first
    for person in ROSTER:
        if person["has_quote_slot"]:
            latest = next(
                (e for e in entries if e.get("user") == person["full_name"] and e.get("comment")),
                None,
            )
            tokens[f"quote_{person['key']}"] = f"“{latest['comment']}”" if latest else ""

    # Leaderboard: rank by count descending, highest-count person's bar is
    # always full width and every other bar is scaled relative to it.
    ranked = sorted(
        ((p, by_user_count.get(p["full_name"], 0)) for p in ROSTER if p["in_leaderboard"]),
        key=lambda pc: pc[1],
        reverse=True,
    )
    max_count = ranked[0][1] if ranked else 0
    leaderboard_bar_widths: dict[int, int] = {}
    for i, (person, count) in enumerate(ranked, start=1):
        tokens[f"lb_name_{i}"] = person["short_name"]
        tokens[f"lb_count_{i}"] = _fmt_num(count)
        fill_shape_id = LEADERBOARD_FILL_SHAPE_IDS[i - 1]
        leaderboard_bar_widths[fill_shape_id] = (
            round(LEADERBOARD_TRACK_WIDTH_EMU * (count / max_count)) if max_count else 0
        )

    # Feedback Rating bars — absolute (score ÷ scale_max), not relative to
    # each other, so a perfect score is always a full bar.
    rating_bar_widths = {}
    for key, shape_id in RATING_BAR_FILL_IDS.items():
        score = rating_scores.get(key)
        frac = min(score / scale_max, 1) if score else 0
        rating_bar_widths[shape_id] = round(RATING_BAR_TRACK_WIDTH_EMU * frac)

    # Ticket Inflow by Area bars — relative to the largest area count.
    area_bar_widths = _bar_widths_by_max(
        {col: area_values[col] for col in AREA_BAR_FILL_IDS}, AREA_BAR_TRACK_WIDTH_EMU
    )
    area_bar_widths = {AREA_BAR_FILL_IDS[col]: w for col, w in area_bar_widths.items()}

    # By Work Type bars — all 10 inflow/outflow bars share one max.
    work_type_bar_widths = _bar_widths_by_max(
        {key: work_type_values[key] for key in WORK_TYPE_BAR_FILL_IDS}, WORK_TYPE_BAR_TRACK_WIDTH_EMU
    )
    work_type_bar_widths = {WORK_TYPE_BAR_FILL_IDS[key]: w for key, w in work_type_bar_widths.items()}

    bar_widths = {
        KPI_SLIDE_INDEX: rating_bar_widths,
        REQUESTS_SLIDE_INDEX: {**area_bar_widths, **work_type_bar_widths},
        LEADERBOARD_SLIDE_INDEX: leaderboard_bar_widths,
    }

    return tokens, bar_widths
